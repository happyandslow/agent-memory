#!/usr/bin/env python3
"""Build the 2-slide KV-transfer deck from PowerPoint_template_sharing.pptx.

Slide 1 — mechanism (the two on-chip topology diagrams).
Slide 2 — measured result (device bandwidth ladder).

Same strategy as 2026-07-20_build.py: keep the template's own slides so the design
is untouched, deep-copy the ones we need, and only replace text / drop placeholders.
Run kv-transfer_figs.py first — it produces the three PNGs this consumes.
"""
import copy, os
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
SRC = os.path.join(HERE, "PowerPoint_template_sharing.pptx")
DST = os.path.join(HERE, "kv-transfer-bandwidth_2page.pptx")
RID = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
RT_IMAGE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"

prs = Presentation(SRC)


def clone_slide(prs, src):
    """Deep-copy a slide (shapes + image relationships) onto the end of the deck."""
    new = prs.slides.add_slide(src.slide_layout)
    for shp in list(new.shapes):
        shp._element.getparent().remove(shp._element)
    remap = {}
    for rel in src.part.rels.values():
        if rel.reltype == RT_IMAGE:
            remap[rel.rId] = new.part.rels._add_relationship(rel.reltype, rel._target)
    for shp in src.shapes:
        el = copy.deepcopy(shp._element)
        for node in el.iter():
            for attr, val in list(node.attrib.items()):
                if attr.startswith('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}') \
                        and val in remap:
                    node.set(attr, remap[val])
        new.shapes._spTree.append(el)
    return new


def by_name(slide, name):
    for shp in slide.shapes:
        if shp.name == name:
            return shp
    raise KeyError(f"{name} not in {[s.name for s in slide.shapes]}")


def set_text(shape, lines, size=None, bold=None):
    """Replace text, reusing the first paragraph/run formatting for every line."""
    if isinstance(lines, str):
        lines = [lines]
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    pPr = copy.deepcopy(p0._p.find(qn('a:pPr'))) if p0._p.find(qn('a:pPr')) is not None else None
    rPr = None
    if p0.runs:
        r0 = p0.runs[0]._r
        rPr = copy.deepcopy(r0.find(qn('a:rPr'))) if r0.find(qn('a:rPr')) is not None else None
    for p in list(tf._txBody.findall(qn('a:p'))):
        tf._txBody.remove(p)
    for line in lines:
        p = tf._txBody.makeelement(qn('a:p'), {})
        if pPr is not None:
            p.append(copy.deepcopy(pPr))
        r = p.makeelement(qn('a:r'), {})
        if rPr is not None:
            r.append(copy.deepcopy(rPr))
        t = r.makeelement(qn('a:t'), {})
        t.text = line
        r.append(t)
        p.append(r)
        tf._txBody.append(p)
    if size is not None or bold is not None:
        for p in tf.paragraphs:
            for r in p.runs:
                if size is not None:
                    r.font.size = Pt(size)
                if bold is not None:
                    r.font.bold = bold


S = prs.slides
keypts, barchart = S[3], S[12]
caption_proto = by_name(barchart, 'Text 1')

s_mech = clone_slide(prs, keypts)
s_res = clone_slide(prs, keypts)


def add_text(slide, left, top, width, height, lines, size, bold=False):
    el = copy.deepcopy(caption_proto._element)
    slide.shapes._spTree.append(el)
    shp = slide.shapes[-1]
    shp.left, shp.top, shp.width, shp.height = (
        Inches(left), Inches(top), Inches(width), Inches(height))
    shp.text_frame.word_wrap = True
    set_text(shp, lines, size=size, bold=bold)
    return shp


# The template's bottom-left footer is a baked IMAGE of the words "Presentation
# title" (shape 'Image 0'), not text — replacing text cannot reach it. Drop it and
# set a real footer line instead.
for sl in (s_mech, s_res):
    for sh in list(sl.shapes):
        if sh.name == 'Image 0':
            sh._element.getparent().remove(sh._element)
    add_text(sl, 1.50, 10.58, 12.00, 0.36,
             "KV cache: prefill → decode transfer  ·  WaferEngine on WSE-3  ·  July 2026", size=11)

# ============================================================ slide 1 — mechanism
set_text(by_name(s_mech, 'Text 0'), "How the KV cache crosses from prefill to decode")
by_name(s_mech, 'Text 1')._element.getparent().remove(by_name(s_mech, 'Text 1')._element)

IMG_H = 6.05
d1_w = IMG_H * 838 / 800          # cropped fig_topo_block.png
d2_w = IMG_H * 875 / 752          # cropped fig_topo_regions.png
GAP = 0.85
x0 = (20.0 - (d1_w + d2_w + GAP)) / 2

add_text(s_mech, x0, 2.46, d1_w, 0.55,
         "1 — inside one prefill block:  funnel → diagonal → transpose        stage A ≈ 7%",
         size=15, bold=True)
add_text(s_mech, x0 + d1_w + GAP, 2.46, d2_w, 0.55,
         "2 — across regions:  store-and-forward north into decode        stage B ≈ 93%",
         size=15, bold=True)

s_mech.shapes.add_picture(os.path.join(HERE, "fig_topo_block.png"),
                          Inches(x0), Inches(3.02), height=Inches(IMG_H))
s_mech.shapes.add_picture(os.path.join(HERE, "fig_topo_regions.png"),
                          Inches(x0 + d1_w + GAP), Inches(3.02), height=Inches(IMG_H))

add_text(s_mech, 1.50, 9.22, 16.90, 1.20, [
    "Prefill and decode hold KV on transposed coordinates — a tile on PE (lx, ly) is expected at (ly, lx). Each row funnels sideways to its diagonal PE, "
    "which then emits down its column; the turn at the diagonal is the transpose. A second, in-tile relayout puts the bytes in decode's order, so decode only receives and unpacks.",
    "That is all of stage A — about 7%. The other 93% is stage B: each tile is handed north PE by PE across the full prefill height and the relay seam, "
    "once per layer per K/V plane, with no overlap between planes.",
], size=13)

# ============================================================== slide 2 — results
set_text(by_name(s_res, 'Text 0'),
         "Measured on real silicon — it scales with geometry, but runs far below the wire")
by_name(s_res, 'Text 1')._element.getparent().remove(by_name(s_res, 'Text 1')._element)

CH_W = 15.4
s_res.shapes.add_picture(os.path.join(HERE, "fig_bw_ladder.png"),
                         Inches((20.0 - CH_W) / 2), Inches(2.62), width=Inches(CH_W))

add_text(s_res, 1.50, 9.22, 16.90, 1.20, [
    "Six configurations end to end on the EPCC CS-3, from a 16×16 toy up to the full 512×512 production geometry — no simulator. Effective bandwidth "
    "= KV bytes moved ÷ (A+B), timed by on-device counters; it grows almost linearly with geometry, reaching 1.803 GB/s.",
    "The ceiling lines are the point: one fabric link carrying a single clean stream measured 3.91 GB/s on the same machine (0.89× spec). The whole "
    "parallel transfer is therefore slower than one clean link and uses ~0.1% of the seam — bound by per-hop latency and by running the planes serially, "
    "not by the fabric. Fixes: coalesce the planes into one shift, enlarge the per-hop payload, cut through, route directly.",
], size=13)

# ------------------------------------------------------- keep only our two slides
keep = [s_mech, s_res]
keep_ids = {id(s) for s in keep}
sldIdLst = prs.slides._sldIdLst
entries = {}
for sldId in list(sldIdLst):
    part = prs.part.rels[sldId.get(RID)].target_part
    entries[id(part.slide)] = sldId
for sid, sldId in entries.items():
    if sid not in keep_ids:
        prs.part.drop_rel(sldId.get(RID))
        sldIdLst.remove(sldId)
for s in keep:
    sldIdLst.append(entries[id(s)])

prs.save(DST)
print("wrote", DST, os.path.getsize(DST), "bytes")
p2 = Presentation(DST)
print("slides:", len(p2.slides._sldIdLst))
for i, s in enumerate(p2.slides):
    t = [sh.text for sh in s.shapes if sh.has_text_frame and sh.text.strip()]
    print(" ", i, "|", (t[0][:78] if t else "(no text)"))
