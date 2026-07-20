#!/usr/bin/env python3
"""Build the weekly deck from PowerPoint_template_sharing.pptx.

Strategy: keep the template's own slides (so the design is untouched) and only
replace their text. Slides that need to appear twice are deep-copied.
"""
import copy, os
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches
from pptx.enum.text import MSO_AUTO_SIZE

SRC = "/home/lexu/agent-memory/projects/WaferEngine-staging/meetings/PowerPoint_template_sharing.pptx"
DST = "/home/lexu/agent-memory/projects/WaferEngine-staging/meetings/2026-07-20.pptx"
RID = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
RT_IMAGE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"

prs = Presentation(SRC)


def clone_slide(prs, src):
    """Deep-copy a slide (shapes + image relationships) onto the end of the deck."""
    new = prs.slides.add_slide(src.slide_layout)
    for shp in list(new.shapes):                    # drop layout placeholders
        shp._element.getparent().remove(shp._element)

    # Re-register the source's image rels. python-pptx assigns its own rIds, so
    # build old -> new and rewrite the references inside the copied XML.
    remap = {}
    for rel in src.part.rels.values():
        if rel.reltype == RT_IMAGE:
            remap[rel.rId] = new.part.rels._add_relationship(rel.reltype, rel._target)

    for shp in src.shapes:
        el = copy.deepcopy(shp._element)
        for node in el.iter():
            for attr, val in list(node.attrib.items()):
                if attr.startswith('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}') \
                        and val in remap:
                    node.set(attr, remap[val])
        new.shapes._spTree.append(el)
    return new


def by_name(slide, name):
    for shp in slide.shapes:
        if shp.name == name:
            return shp
    raise KeyError(f"{name} not in {[s.name for s in slide.shapes]}")


def set_text(shape, lines):
    """Replace text, reusing the first paragraph/run formatting for every line."""
    if isinstance(lines, str):
        lines = [lines]
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    pPr = copy.deepcopy(p0._p.find(qn('a:pPr'))) if p0._p.find(qn('a:pPr')) is not None else None
    rPr = None
    if p0.runs:
        r0 = p0.runs[0]._r
        rPr = copy.deepcopy(r0.find(qn('a:rPr'))) if r0.find(qn('a:rPr')) is not None else None

    for p in list(tf._txBody.findall(qn('a:p'))):   # clear
        tf._txBody.remove(p)
    for line in lines:
        p = tf._txBody.makeelement(qn('a:p'), {})
        if pPr is not None:
            p.append(copy.deepcopy(pPr))
        r = p.makeelement(qn('a:r'), {})
        if rPr is not None:
            r.append(copy.deepcopy(rPr))
        t = r.makeelement(qn('a:t'), {})
        t.text = line
        r.append(t)
        p.append(r)
        tf._txBody.append(p)


# ---------------------------------------------------------------- clone first
S = prs.slides
title, keypts, twocol, table, barchart = S[0], S[3], S[4], S[11], S[12]
keypts2 = clone_slide(prs, keypts)
diag_p = clone_slide(prs, keypts)
diag_d = clone_slide(prs, keypts)
table2 = clone_slide(prs, table)

# a caption shape to reuse under the two data tables
caption_proto = by_name(barchart, 'Text 1')

# ------------------------------------------------------------------- 1. title
set_text(by_name(title, 'Text 0'), "WaferEngine — Weekly Update")
set_text(by_name(title, 'Text 1'),
         "M0 reuse foundation: PE-internal KV reuse (S6a) — prefill warm-start lands, "
         "and both kernels are measured on real WSE-3 for the first time")
set_text(by_name(title, 'Text 2'), "Le Xu")
set_text(by_name(title, 'Text 4'), "20 July 2026")

for _s in (keypts, keypts2):
    _b = by_name(_s, 'Text 1')
    _b.left, _b.width = Inches(1.47), Inches(16.90)
_b = by_name(keypts2, 'Text 1'); _b.top, _b.height = Inches(3.30), Inches(6.20)

# master/layout footer still says "Presentation title"
for _cand in list(prs.slide_layouts) + [prs.slide_master]:
    for _sh in _cand.shapes:
        if _sh.has_text_frame and 'Presentation title' in _sh.text_frame.text:
            set_text(_sh, "WaferEngine — Weekly Update  ·  20 July 2026")

# ----------------------------------------------------------------- 2. roadmap
set_text(by_name(keypts, 'Text 0'), "Roadmap — and where we are")
set_text(by_name(keypts, 'Text 1'), [
    "North star — tiered KV cache management (in-PE → idle-PE → off-chip) plus a policy "
    "control plane, under a hard on-chip SRAM budget.",
    "",
    "M0  Reuse foundation — preserve and key-retrieve KV across consecutive requests.   ◀  WE ARE HERE",
    "M1  Single-card in-place reuse — skip prefill on a prefix hit; requests-per-bank capacity curve.",
    "M2  force-decode vs ship-to-prefill — an A/B cost model, validated on a multi-turn trace.",
    "M3  Idle-PE second tier — park KV on free PEs and reload it; quantify the real cost.",
    "M4  Off-chip retained pool + eviction — and lift the ~512-token single-pass prefill cap.",
    "M5  Unified control plane — one min-waste policy across all tiers, ablatable and explainable.",
    "M6  Scenario evaluation — an agentic workload, plus an area/power-equivalent GPU comparison.",
])

# ---------------------------------------------------------------- 3. progress
set_text(by_name(keypts2, 'Text 0'), "This week — M0 / S6a")
set_text(by_name(keypts2, 'Text 1'), [
    "S6a-decode  ✓  done earlier and re-confirmed: decode keeps KV across requests instead of "
    "discarding it; output matches the reference to the noise floor.",
    "S6a-prefill ✓  the week's main deliverable: a request can now start from a resident prefix "
    "instead of from token 0. Implemented, and verified at full scale.",
    "",
    "Along the way we found a hardware-level trap: a fabric transfer whose length is odd never "
    "signals completion on WSE-3 — it just hangs, with no compile error. Reproduced in a small "
    "standalone test, then fixed by padding the per-request metadata to an even width.",
    "Three correctness bugs in the warm-start path were found and fixed — all three would have "
    "produced plausible-looking but wrong output rather than an obvious failure.",
    "",
    "First real-scale device runs: 524,288 PEs, 28 layers, 8,192 tokens. Every reuse round is "
    "byte-identical to the cold round, so the saving is real, not an accuracy trade.",
    "Performance measured for both kernels — next two slides.",
    "",
    "Open: a longer-sequence (16,384-token) sweep is running now; code is not committed yet, "
    "pending review.",
])

# ------------------------------- 4a/4b. implementation, as diagrams not prose
for _sl, _title, _img in (
        (diag_p, "What changed — prefill: start from chunk k, not chunk 0", "diag_prefill.png"),
        (diag_d, "What changed — decode: retain the KV instead of rewinding", "diag_decode.png")):
    set_text(by_name(_sl, 'Text 0'), _title)
    _body = by_name(_sl, 'Text 1')                     # drop the text placeholder
    _body._element.getparent().remove(_body._element)
    _sl.shapes.add_picture(os.path.join(os.path.dirname(os.path.abspath(__file__)), _img),
                           Inches(1.47), Inches(2.55), width=Inches(17.06))

# ------------------------------------------------------------ 5. prefill perf
def fill_table(slide, title_text, header, rows, caption):
    set_text(by_name(slide, 'Text 0'), title_text)
    for i, h in enumerate(header):
        shp = by_name(slide, f'Text {2 + i * 2}')
        # NOTE: rendering these headers with LibreOffice clips the last glyph
        # (ITEM -> ITEN). The untouched template does exactly the same, so it is
        # a missing-brand-font substitution artifact of that renderer, not the
        # file. Do not "fix" it by forcing width/wrap -- that breaks real
        # PowerPoint, where it renders correctly.
        set_text(shp, h)
    for r, row in enumerate(rows):
        base = 12 + r * 10
        for c, val in enumerate(row):
            set_text(by_name(slide, f'Text {base + c * 2}'), val)
    cap = copy.deepcopy(caption_proto._element)
    slide.shapes._spTree.append(cap)
    shp = slide.shapes[-1]
    shp.left, shp.top = Inches(1.50), Inches(9.20)
    shp.width, shp.height = Inches(16.90), Inches(1.10)
    shp.text_frame.word_wrap = True
    set_text(shp, caption)


fill_table(
    table,
    "Prefill — reusing half the prompt saves less than a quarter of the time",
    ["PREFIX", "SHARE", "LATENCY", "SAVING", "THROUGHPUT"],
    [["none",      "0 %",  "1001.5 ms", "—",       "8,180 tok/s"],
     ["8 chunks",  "25 %", "924.1 ms",  "−7.7 %",  "8,865 tok/s"],
     ["16 chunks", "50 %", "773.3 ms",  "−22.8 %", "10,594 tok/s"],
     ["24 chunks", "75 %", "549.2 ms",  "−45.2 %", "14,916 tok/s"]],
    "8,192-token prompt on 524,288 PEs; pure on-device request latency, host work excluded. "
    "All reuse rounds byte-identical to the cold round.   —   Why sub-linear: a chunk's cost "
    "grows with how much text precedes it (the last chunk costs about 7× the first), so a "
    "reused prefix is always the cheapest part of the request. Planning that assumes "
    "“half the prompt cached = half the time saved” will overestimate the win.",
)

# ------------------------------------------------------------- 6. decode perf
fill_table(
    table2,
    "Decode — the win is skipping whole steps, not making steps cheaper",
    ["MEASURED", "ARM", "COST", "CHANGE", "MEANING"],
    [["Second round",        "no reuse",          "262.9 M cycles",   "—",
      "redoes 512 discarded steps"],
     ["Second round",        "retain",            "127.7 M cycles",   "−51.4 %",
      "runs only the 256 new steps"],
     ["Total decode work",   "both rounds",       "390.6 → 255.4 M", "−34.6 %",
      "as predicted: 768 → 512 steps"],
     ["Control: equal work", "same steps both", "126.2 vs 126.2 M", "+0.02 %",
      "a retained step is not cheaper"]],
    "524,288 PEs; cycles at 1.1 GHz; timer starts after KV load, so this is compute only. "
    "The control row is the point: retain changes the cost of a decode step by 0.02 %. "
    "Its entire value is that steps already taken do not have to be taken again — which is why "
    "the comparison has to make the no-reuse arm redo the work it threw away.",
)

# -------------------------------------------------- keep only our six slides
keep = [title, keypts, keypts2, diag_p, diag_d, table, table2]
keep_ids = {id(s) for s in keep}
sldIdLst = prs.slides._sldIdLst
entries = {}
for sldId in list(sldIdLst):
    part = prs.part.rels[sldId.get(RID)].target_part
    entries[id(part.slide)] = sldId
for sid, sldId in entries.items():
    if sid not in keep_ids:
        prs.part.drop_rel(sldId.get(RID))
        sldIdLst.remove(sldId)
for s in keep:                                  # reorder
    sldIdLst.append(entries[id(s)])

prs.save(DST)
print("wrote", DST, os.path.getsize(DST), "bytes")
p2 = Presentation(DST)
print("slides:", len(p2.slides.__iter__.__self__._sldIdLst))
for i, s in enumerate(p2.slides):
    t = [sh.text for sh in s.shapes if sh.has_text_frame and sh.text.strip()]
    print(i, "|", (t[0][:70] if t else "(no text)"))
