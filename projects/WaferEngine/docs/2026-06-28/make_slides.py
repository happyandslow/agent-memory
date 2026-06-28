"""Generate a concise slide deck (PPTX) summarizing the per-PE resource analysis
of qwen3-1.7B on the Cerebras WSE-3. For non-expert viewers.

  python projects/WaferEngine/docs/2026-06-28/make_slides.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
# Source plots live in the WaferEngine repo (branch lexu/pe-mem-breakdown); the
# generated deck is committed here in agent-memory.
RES = os.environ.get(
    "PE_MEM_RESULTS",
    "/home/lexu/WaferEngine/tools/pe_mem_breakdown/results")
OUT = os.path.join(HERE, "wse_per_pe_resource_analysis.pptx")

NAVY = RGBColor(0x1F, 0x3A, 0x5F)
BLUE = RGBColor(0x1F, 0x77, 0xB4)
GREY = RGBColor(0x55, 0x55, 0x55)
DARK = RGBColor(0x22, 0x22, 0x22)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _tb(slide, l, t, w, h):
    box = slide.shapes.add_textbox(l, t, w, h)
    box.text_frame.word_wrap = True
    return box.text_frame


def title_slide(title, subtitle):
    s = prs.slides.add_slide(BLANK)
    bar = s.shapes.add_shape(1, 0, Inches(2.6), SW, Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    tf = _tb(s, Inches(0.8), Inches(2.8), Inches(11.7), Inches(2.2))
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = NAVY
    p2 = tf.add_paragraph(); p2.text = subtitle
    p2.font.size = Pt(20); p2.font.color.rgb = GREY
    return s


def content_slide(title, bullets, img=None, img_caption=None, img2=None):
    s = prs.slides.add_slide(BLANK)
    bar = s.shapes.add_shape(1, 0, 0, SW, Inches(0.95))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    ttf = _tb(s, Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.75))
    tp = ttf.paragraphs[0]; tp.text = title
    tp.font.size = Pt(26); tp.font.bold = True; tp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    text_w = Inches(6.0) if img else Inches(12.5)
    btf = _tb(s, Inches(0.5), Inches(1.2), text_w, Inches(5.9))
    for i, b in enumerate(bullets):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        lvl = 0
        if b.startswith("  "):
            lvl = 1; b = b.strip()
        p.text = ("• " if lvl == 0 else "– ") + b
        p.level = lvl
        p.font.size = Pt(18 if lvl == 0 else 15)
        p.font.color.rgb = DARK if lvl == 0 else GREY
        p.space_after = Pt(7)

    if img:
        imgs = [img] if not img2 else [img, img2]
        x = Inches(6.7)
        avail_w = Inches(6.3)
        if len(imgs) == 1:
            pic = s.shapes.add_picture(imgs[0], x, Inches(1.25), height=Inches(5.4))
            if pic.width > avail_w:
                pic.height = int(pic.height * avail_w / pic.width); pic.width = avail_w
            pic.left = x + int((avail_w - pic.width) / 2)
        else:
            for k, im in enumerate(imgs):
                pic = s.shapes.add_picture(im, x, Inches(1.25 + k * 2.85), height=Inches(2.7))
                if pic.width > avail_w:
                    pic.height = int(pic.height * avail_w / pic.width); pic.width = avail_w
                pic.left = x + int((avail_w - pic.width) / 2)
        if img_caption:
            cf = _tb(s, x, Inches(6.75), avail_w, Inches(0.5))
            cp = cf.paragraphs[0]; cp.text = img_caption
            cp.font.size = Pt(11); cp.font.italic = True; cp.font.color.rgb = GREY
            cp.alignment = PP_ALIGN.CENTER
    return s


def image_slide(title, img, bullets):
    s = prs.slides.add_slide(BLANK)
    bar = s.shapes.add_shape(1, 0, 0, SW, Inches(0.95))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    ttf = _tb(s, Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.75))
    tp = ttf.paragraphs[0]; tp.text = title
    tp.font.size = Pt(26); tp.font.bold = True; tp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    pic = s.shapes.add_picture(img, Inches(0.4), Inches(1.1), height=Inches(5.0))
    if pic.width > Inches(8.6):
        pic.height = int(pic.height * Inches(8.6) / pic.width); pic.width = Inches(8.6)
    btf = _tb(s, Inches(9.2), Inches(1.2), Inches(3.9), Inches(5.9))
    for i, b in enumerate(bullets):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.text = "• " + b; p.font.size = Pt(15); p.font.color.rgb = DARK
        p.space_after = Pt(8)
    return s


D = os.path.join(RES, "decode")
P = os.path.join(RES, "prefill")

# 1 ----------------------------------------------------------------- title
title_slide("Serving an LLM on a Wafer: a per-PE resource X-ray",
            "qwen3-1.7B on a real Cerebras WSE-3 (CS-3)  ·  memory & fabric-color "
            "analysis  ·  2026-06-28")

# 2 ------------------------------------------------------------ why it matters
content_slide(
    "Why this is hard: no shared memory",
    ["A GPU has one big shared memory (HBM). A Cerebras wafer is ~900,000 tiny "
     "cores (PEs), each with its own tiny 48 KB of memory — nothing shared.",
     "To run an LLM you must REPLICATE the program on every core and SPLIT the "
     "weights across them.",
     "So each core has very little room left for the KV cache (the running "
     "context) — which caps the sequence length.",
     "Two scarce per-core resources: (1) 48 KB SRAM, and (2) 24 'colors' — the "
     "fabric communication channels.",
     "Question: what does ONE core actually hold and use? We built a tool to "
     "find out, and ran it on real silicon."])

# 3 -------------------------------------------------------------- method
content_slide(
    "Method: X-ray every PE on real hardware",
    ["Compile each kernel (decode & prefill) on a real CS-3, for the full "
     "762x1172 fabric.",
     "Read each PE's memory map (cs-readelf) + the compiled binaries, and "
     "classify every byte: code / weights / KV cache / activations / system / free.",
     "Covered all ~329,000 placed PEs — but the result is really PER-ROLE "
     "(every PE of the same role is identical), so it's a fast static analysis.",
     "Also analyzed fabric-COLOR usage per microstep from the kernel source.",
     "Tool: tools/pe_mem_breakdown/  (28 tests, reconciles to the byte)."])

# 4 -------------------------------------------------- memory breakdown
image_slide(
    "Finding 1: code eats half of every core",
    os.path.join(D, "decode_stacked.png"),
    ["Each compute PE: ~22 KB CODE + ~11 KB weights.",
     "CODE = the full transformer-layer program, copied onto every one of "
     "~65,000 compute cores.",
     "KV cache is tiny (<0.3 KB) — it gets only the leftover ~11 KB.",
     "ht_head/ht_tail cores are weight-bound (embedding / output table)."])

# 5 ------------------------------------------------ spatial placement
image_slide(
    "Where the blocks live on the wafer",
    os.path.join(D, "decode_placement_map.png"),
    ["The kernel is laid out as functional blocks on the 2-D grid.",
     "Big area = the 2x2 compute blocks (the transformer layers).",
     "Thin left band = embedding (ht_head) + output/sampling (ht_tail).",
     "We can map any metric (code, weights, KV, free...) to its (x,y) "
     "coordinate as a heatmap."])

# 6 ------------------------------------------------ key memory insight
content_slide(
    "Finding 1 — what it means",
    ["The 'disaggregated-memory tax': you pay for the whole program ~65,000 "
     "times, before a single weight or KV byte.",
     "Cores of the same role are memory-IDENTICAL on silicon (verified over all "
     "329K PEs) — no surprises, no per-core variation.",
     "Busy cores run ~77-82% full. KV is the squeezed residual.",
     "Consequence: max sequence length is limited by CODE + WEIGHTS, not by a "
     "'KV cache size' choice. To go longer you must shrink code/weights."],
    img=os.path.join(D, "heatmaps", "decode_heatmaps_all.png"),
    img_caption="Per-PE memory by fabric coordinate (8 fields)")

# 7 ------------------------------------------------ max seq len
content_slide(
    "Finding 2: max context = 22,784 tokens (44x the default)",
    ["qwen3-1.7B is light, so the leftover ~11 KB is real headroom for KV.",
     "We swept the sequence length on real hardware (compile + place):",
     "  PASS up to 22,784 tokens; FAILS (out of memory) at 23,040.",
     "  That's 44.5x the shipped default of 512.",
     "KV costs ~0.44 bytes/token per core; the ~11 KB free fills until the last "
     "~1 KB (needed for placement) is gone.",
     "Caveat: this is the compile/fit ceiling, for one request (batch size 1)."])

# 8 ------------------------------------------------ colors
image_slide(
    "Finding 3: 'colors' are the other scarce resource",
    os.path.join(D, "decode_color_usage.png"),
    ["24 fabric 'colors' = communication channels. We mapped which are used "
     "in each microstep.",
     "DECODE uses only 8/24 on a compute core...",
     "...because it REUSES 5 channels, repainting their routes 6x per layer "
     "(Y / X / band = the reduce direction).",
     "Local steps (RoPE, residuals) use no fabric at all."])

# 9 ------------------------------------------ decode vs prefill colors
content_slide(
    "The decode <-> prefill trade-off",
    ["DECODE = color-frugal, reconfiguration-heavy: 8/24 colors; reuses 5 and "
     "repaints their routes every step.",
     "PREFILL = color-hungry, reconfiguration-light: 17/24 colors.",
     "  Its matmuls use a systolic 'MeshGEMM' with 6 dedicated, statically-"
     "routed channels (never repainted mid-matmul).",
     "So the two kernels trade the SAME two resources differently: decode spends "
     "route-repaint time to save colors; prefill spends colors to save time.",
     "Both fit comfortably — neither is color-bound today."],
    img=os.path.join(P, "prefill_color_usage.png"),
    img_caption="Prefill: 17/24 colors (MeshGEMM-heavy)")

# 10 ---------------------------------------------- takeaways + traceability
content_slide(
    "Takeaways & where to find everything",
    ["Two scarce per-core resources on the wafer: 48 KB SRAM and 24 colors.",
     "SRAM: code (~half) is the dominant, replicated cost; KV is the residual; "
     "max context = 22,784 tokens (measured).",
     "Colors: decode 8/24 (reconfig-heavy), prefill 17/24 (MeshGEMM-heavy); "
     "neither is color-bound.",
     "All analysis is per-role static -> cheap and reproducible.",
     "Traceable: branch  lexu/pe-mem-breakdown   (WaferEngine repo).",
     "  tool + plots: tools/pe_mem_breakdown/   ·   results/{decode,prefill}/",
     "  key commits: 33ecf27 (heatmaps), 023cf1d (seq-len probe), "
     "0869200 (color map)."])

prs.save(OUT)
print("wrote", OUT, f"({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
